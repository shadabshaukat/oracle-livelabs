from __future__ import annotations

import os
import re
import logging
import subprocess
import zipfile
import shutil
from dataclasses import dataclass
from typing import List, Tuple, Iterable
import csv
import json
from tempfile import NamedTemporaryFile

from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
from pypdf import PdfReader
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from .config import settings

logger = logging.getLogger(__name__)

# Safety guardrails are controlled via environment-backed settings.


# Embedded-image OCR fallback guard to avoid repeated per-image warnings when OCR backend
# is not available (e.g., tesseract missing on host).
_EMBEDDED_IMAGE_OCR_DISABLED = False
_EMBEDDED_IMAGE_OCR_WARNED = False


def _disable_embedded_image_ocr_once(reason: str) -> None:
    global _EMBEDDED_IMAGE_OCR_DISABLED, _EMBEDDED_IMAGE_OCR_WARNED
    _EMBEDDED_IMAGE_OCR_DISABLED = True
    if not _EMBEDDED_IMAGE_OCR_WARNED:
        logger.warning("Embedded-image OCR disabled for this process: %s", reason)
        _EMBEDDED_IMAGE_OCR_WARNED = True


# Prefer keeping paragraph boundaries; avoid collapsing all newlines into spaces
PARA_SPLIT_RE = re.compile(r"(?:\r?\n){2,}")
SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")
UPPER_HEADING_RE = re.compile(r"^[A-Z0-9][A-Z0-9 \-:]{2,}$")
NUMBERED_HEADING_RE = re.compile(r"^(?:[IVXLCDM]+\.|\d+(?:\.\d+)*\.|[A-Z]\.)\s+.+")
PAGE_FOOTER_RE = re.compile(r"^\s*page\s+\d+(?:\s+of\s+\d+)?\s*$", re.I)


@dataclass
class ChunkParams:
    chunk_size: int = 2500
    chunk_overlap: int = 250
    strategy: str = "recursive"  # recursive | sentence_pack
    # Optional custom separator order for recursive splitting
    separators: tuple[str, ...] = ("\n\n", "\n", ". ", " ", "")


def _normalize_whitespace_preserve_paragraphs(text: str) -> str:
    """Normalize whitespace but preserve blank lines as paragraph boundaries."""
    # Remove NUL bytes that can break database inserts
    text = text.replace("\x00", "")
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse more than two consecutive newlines to exactly two
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Normalize spaces within lines
    lines = []
    for ln in text.split("\n"):
        ln = re.sub(r"\s+", " ", ln).strip()
        lines.append(ln)
    text = "\n".join(lines)
    # Restore paragraph boundaries
    text = re.sub(r"(\n\s*){3,}", "\n\n", text)
    return text.strip()


def _fix_hyphenation(text: str) -> str:
    """Fix common PDF hyphenation like 'exam-\nple' -> 'example'."""
    # Join words broken by hyphen at line end
    text = re.sub(r"-\n(?=\w)", "", text)
    # Remove lone hyphens surrounded by newlines
    text = re.sub(r"\n-\n", "\n", text)
    # Replace single newlines inside paragraphs with spaces (but keep double newlines)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"(?<!\n)\n(?!\n)", " ", text)
    return text


def _insert_heading_boundaries(text: str) -> str:
    """Insert extra blank lines around detected headings to improve chunk boundaries."""
    out_lines: List[str] = []
    for ln in text.split("\n"):
        if UPPER_HEADING_RE.match(ln) or NUMBERED_HEADING_RE.match(ln):
            if out_lines and out_lines[-1] != "":
                out_lines.append("")
            out_lines.append(ln)
            out_lines.append("")
        else:
            out_lines.append(ln)
    return "\n".join(out_lines)


def _remove_common_headers_footers(pages: List[str]) -> List[str]:
    """Heuristic removal of repeating headers/footers across pages."""
    if not pages or len(pages) < 3:
        return pages
    # Collect first and last non-empty lines per page
    first_lines: List[str] = []
    last_lines: List[str] = []
    for p in pages:
        ls = [l.strip() for l in p.split("\n") if l.strip()]
        if not ls:
            first_lines.append("")
            last_lines.append("")
            continue
        first_lines.append(ls[0])
        last_lines.append(ls[-1])
    def most_common(cand: List[str]) -> str:
        from collections import Counter
        c = Counter([x for x in cand if x])
        return c.most_common(1)[0][0] if c else ""
    first_common = most_common(first_lines)
    last_common = most_common(last_lines)
    cleaned_pages: List[str] = []
    for p in pages:
        ls = p.split("\n")
        if first_common:
            ls = ls[1:] if ls and ls[0].strip() == first_common else ls
        if last_common:
            if ls and ls[-1].strip() == last_common:
                ls = ls[:-1]
        # Remove generic page footers like "Page X of Y"
        ls = [l for l in ls if not PAGE_FOOTER_RE.match(l.strip())]
        cleaned_pages.append("\n".join(ls))
    return cleaned_pages


def read_text_from_file(path: str) -> Tuple[str, str]:
    """
    Return (text, source_type) from a supported file.
    source_type: pdf|html|txt|docx|doc|xml|csv|md|json
    """
    ext = os.path.splitext(path)[1].lower()
    if not os.path.exists(path):
        logger.error("File does not exist for extraction: %s", path)
        return "", ext.lstrip(".") or "txt"
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}:
        return "", "image"
    if ext == ".pdf":
        try:
            return extract_text_from_pdf(path), "pdf"
        except Exception as exc:
            logger.exception("Failed to extract PDF text from %s: %s", path, exc)
            return "", "pdf"
    if ext in {".html", ".htm"}:
        try:
            return extract_text_from_html(path), "html"
        except Exception as exc:
            logger.exception("Failed to extract HTML text from %s: %s", path, exc)
            return "", "html"
    if ext in {".docx", ".doc"}:
        try:
            if ext == ".doc":
                return extract_text_from_doc(path), "doc"
            return extract_text_from_docx(path), "docx"
        except Exception as exc:
            logger.exception("Failed to extract DOC/DOCX text from %s: %s", path, exc)
            return "", ext.lstrip(".")
    if ext in {".txt", ""}:
        try:
            return extract_text_from_txt(path), "txt"
        except Exception as exc:
            logger.exception("Failed to extract TXT text from %s: %s", path, exc)
            return "", "txt"
    if ext == ".xml":
        try:
            return extract_text_from_xml(path), "xml"
        except Exception as exc:
            logger.exception("Failed to extract XML text from %s: %s", path, exc)
            return "", "xml"
    if ext == ".csv":
        try:
            return extract_text_from_csv(path), "csv"
        except Exception as exc:
            logger.exception("Failed to extract CSV text from %s: %s", path, exc)
            return "", "csv"
    if ext == ".md":
        try:
            return extract_text_from_md(path), "md"
        except Exception as exc:
            logger.exception("Failed to extract Markdown text from %s: %s", path, exc)
            return "", "md"
    if ext == ".json":
        try:
            return extract_text_from_json(path), "json"
        except Exception as exc:
            logger.exception("Failed to extract JSON text from %s: %s", path, exc)
            return "", "json"
    if ext == ".pptx":
        try:
            return extract_text_from_pptx(path), "pptx"
        except Exception as exc:
            logger.exception("Failed to extract PPTX text from %s: %s", path, exc)
            return "", "pptx"
    if ext in {".xlsx", ".xls"}:
        try:
            return extract_text_from_excel(path), ext.lstrip(".")
        except Exception as exc:
            logger.exception("Failed to extract spreadsheet text from %s: %s", path, exc)
            return "", ext.lstrip(".")
    # Fallback: read as text if possible
    try:
        return extract_text_from_txt(path), ext.lstrip('.') or 'txt'
    except Exception:
        raise ValueError(f"Unsupported file type: {ext}")


def extract_text_from_pdf(path: str) -> str:
    """Robust PDF extraction.
    Order of preference:
      1) PyMuPDF (if enabled): page.get_text("text"); remove common headers/footers; hyphenation fix; preserve paragraphs
      2) pypdf: page.extract_text(); hyphenation fix; preserve paragraphs
      3) pdfplumber fallback for table/figure-heavy PDFs when pypdf output is sparse
    """
    text_pymupdf = ""
    # Optional: use PyMuPDF if enabled and available for better extraction
    if getattr(settings, "use_pymupdf", False):
        try:
            import fitz  # PyMuPDF
            pages_raw: List[str] = []
            with fitz.open(path) as doc:
                for page in doc:
                    # Use textual extraction; "text" preserves reading order better than "blocks" in many docs
                    t = page.get_text("text") or ""
                    pages_raw.append(t)
            # Remove common headers/footers
            pages_clean = _remove_common_headers_footers(pages_raw)
            text_pymupdf = "\n\n".join(pages_clean)
            text_pymupdf = _fix_hyphenation(text_pymupdf)
            text_pymupdf = _normalize_whitespace_preserve_paragraphs(text_pymupdf)
            # Insert heading boundaries to help chunking
            text_pymupdf = _insert_heading_boundaries(text_pymupdf)
        except Exception:
            # Fall back to other extractors if PyMuPDF is not available or fails
            text_pymupdf = ""

    # pypdf extraction
    reader = PdfReader(path)
    texts_pypdf: List[str] = []
    try:
        for page in reader.pages:
            txt = page.extract_text() or ""
            texts_pypdf.append(txt)
            # Attempt to include tables when available
            try:
                tables = page.extract_tables()
                if tables:
                    texts_pypdf.append(_tables_to_text(tables))
            except Exception:
                pass
    except Exception:
        texts_pypdf = []
    text_pypdf = "\n\n".join(texts_pypdf)
    text_pypdf = _fix_hyphenation(text_pypdf)
    text_pypdf = _normalize_whitespace_preserve_paragraphs(text_pypdf)
    text_pypdf = _insert_heading_boundaries(text_pypdf)

    # Decide if we should try pdfplumber fallback (very sparse output or extremely short)
    needs_fallback = len(text_pypdf.strip()) < 200 or text_pypdf.count("\n") < max(2, len(texts_pypdf) // 4)
    text_plumb = ""
    if needs_fallback:
        try:
            import pdfplumber  # type: ignore
            with pdfplumber.open(path) as pdf:
                pages_text = []
                for page in pdf.pages:
                    # Tolerances can help capture columns/tables better
                    t = page.extract_text(x_tolerance=1, y_tolerance=1) or ""
                    tables = page.extract_tables() or []
                    if tables:
                        t = "\n\n".join([t, _tables_to_text(tables)]) if t else _tables_to_text(tables)
                    pages_text.append(t)
            text_plumb = "\n\n".join(pages_text)
            text_plumb = _fix_hyphenation(text_plumb)
            text_plumb = _normalize_whitespace_preserve_paragraphs(text_plumb)
            text_plumb = _insert_heading_boundaries(text_plumb)
        except Exception:
            # If pdfplumber fails, keep pypdf output
            text_plumb = ""

    # Choose best non-OCR text result
    candidates = [t for t in [text_pymupdf, text_plumb, text_pypdf] if t and t.strip()]
    best_text = max(candidates, key=lambda t: len(t.strip()), default="")

    # OCR is a fallback for scanned/sparse PDFs. Appending OCR to already-good
    # native text duplicates content and lowers retrieval quality.
    should_try_ocr = (
        settings.ocr_pdf_enabled
        and len((best_text or "").strip()) < max(1, int(settings.pdf_ocr_fallback_min_chars))
    )
    if should_try_ocr:
        try:
            ocr_text = _ocr_pdf_pages(path)
        except Exception as exc:
            logger.warning("PDF OCR fallback failed for %s: %s", path, exc)
            ocr_text = ""
        if ocr_text:
            combined = "\n\n".join([t for t in [best_text, ocr_text] if t and t.strip()])
            combined = _normalize_whitespace_preserve_paragraphs(combined)
            return combined

    return best_text or text_pypdf


def _ocr_pdf_pages(path: str) -> str:
    """OCR PDF pages by rasterizing with PyMuPDF and running OCR per page."""
    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        logger.warning("PyMuPDF not available for PDF OCR: %s", exc)
        return ""
    try:
        from .vision_embeddings import ocr_image_text, OcrUnavailable
    except Exception as exc:
        logger.warning("OCR pipeline unavailable: %s", exc)
        return ""
    texts: List[str] = []
    with fitz.open(path) as doc:
        for page_num, page in enumerate(doc, start=1):
            try:
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                with NamedTemporaryFile(delete=False, suffix=f"_p{page_num}.png") as tmp:
                    pix.save(tmp.name)
                    tmp_path = tmp.name
                try:
                    text = ocr_image_text(tmp_path)
                    if text:
                        texts.append(text)
                finally:
                    try:
                        os.remove(tmp_path)
                    except OSError:
                        pass
            except OcrUnavailable as exc:
                logger.warning("OCR unavailable for PDF page %s: %s", page_num, exc)
                break
            except Exception as exc:
                logger.warning("Failed OCR on PDF page %s: %s", page_num, exc)
                continue
    return "\n\n".join([t for t in texts if t and t.strip()])


def _ocr_text_from_image_bytes(image_bytes: bytes, *, suffix: str = ".png") -> str:
    if _EMBEDDED_IMAGE_OCR_DISABLED:
        return ""
    if not image_bytes or not settings.ocr_enabled:
        return ""
    try:
        from .vision_embeddings import ocr_image_text, OcrUnavailable
    except Exception as exc:
        _disable_embedded_image_ocr_once(str(exc))
        return ""
    with NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(image_bytes)
        tmp_path = tmp.name
    try:
        return ocr_image_text(tmp_path) or ""
    except OcrUnavailable as exc:
        msg = str(exc)
        lowered = msg.lower()
        if (
            "tesseract is not installed" in lowered
            or "not in your path" in lowered
            or "pytesseract is not installed" in lowered
        ):
            _disable_embedded_image_ocr_once(msg)
        else:
            logger.warning("OCR unavailable for embedded image: %s", exc)
        return ""
    except Exception as exc:
        logger.warning("Embedded image OCR failed: %s", exc)
        return ""
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass


def _extract_embedded_image_ocr_from_zip(path: str, media_prefix: str) -> str:
    if not settings.ocr_enabled or not zipfile.is_zipfile(path):
        return ""
    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"}
    texts: List[str] = []
    processed = 0
    total_bytes = 0
    try:
        with zipfile.ZipFile(path) as zf:
            for name in zf.namelist():
                if _EMBEDDED_IMAGE_OCR_DISABLED:
                    break
                if not name.startswith(media_prefix):
                    continue
                ext = os.path.splitext(name)[1].lower()
                if ext not in image_exts:
                    continue
                if processed >= settings.embedded_ocr_max_images:
                    logger.info("Embedded-image OCR capped at %s images for %s", settings.embedded_ocr_max_images, path)
                    break
                try:
                    img_bytes = zf.read(name)
                except Exception:
                    continue
                if not img_bytes:
                    continue
                total_bytes += len(img_bytes)
                if total_bytes > settings.embedded_ocr_max_total_bytes:
                    logger.info(
                        "Embedded-image OCR byte cap reached for %s (>%s bytes)",
                        path,
                        settings.embedded_ocr_max_total_bytes,
                    )
                    break
                ocr_txt = _ocr_text_from_image_bytes(img_bytes, suffix=ext or ".png")
                processed += 1
                if ocr_txt:
                    texts.append(ocr_txt)
    except Exception as exc:
        logger.warning("Failed embedded-image OCR scan for %s: %s", path, exc)
        return ""
    return "\n\n".join([t for t in texts if t and t.strip()])


def extract_text_from_html(path: str) -> str:
    with open(path, "rb") as f:
        data = f.read()
    soup = BeautifulSoup(data, "html.parser")
    # Remove nav-like elements
    for tag in soup(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_xml(path: str) -> str:
    with open(path, "rb") as f:
        data = f.read()
    soup = BeautifulSoup(data, "xml")
    if soup:
        # Prefer element text but include attribute values for structured XML
        parts: List[str] = []
        for el in soup.find_all(True):
            if el.name:
                attrs = " ".join(f"{k}={v}" for k, v in el.attrs.items())
                if attrs:
                    parts.append(attrs)
            if el.string and el.string.strip():
                parts.append(el.string.strip())
        text = "\n".join(parts) or soup.get_text(separator="\n", strip=True)
    else:
        text = ""
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_csv(path: str) -> str:
    parts: List[str] = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        sample = f.read(4096)
        f.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel
        reader = csv.reader(f, dialect)
        for row in reader:
            cleaned = [cell.strip() for cell in row]
            if any(cleaned):
                parts.append("\t".join(cleaned))
    text = "\n".join(parts)
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_md(path: str) -> str:
    # Light-weight: treat as plain text (could add markdown parsing if needed)
    return extract_text_from_txt(path)


def extract_text_from_json(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            data = json.load(f)
        # Convert JSON to a flat text string
        def _flatten(obj) -> List[str]:
            out: List[str] = []
            if isinstance(obj, dict):
                for k, v in obj.items():
                    out.append(str(k))
                    out.extend(_flatten(v))
            elif isinstance(obj, list):
                for it in obj:
                    out.extend(_flatten(it))
            else:
                out.append(str(obj))
            return out
        parts = _flatten(data)
        text = "\n".join(s.strip() for s in parts if s and isinstance(s, str))
        text = _normalize_whitespace_preserve_paragraphs(text)
        return text
    except Exception:
        # Fall back to raw text if not valid JSON
        return extract_text_from_txt(path)


def _looks_like_text(text: str) -> bool:
    if not text:
        return False
    sample = text[:4000]
    if not sample.strip():
        return False
    printable = sum(1 for ch in sample if ch.isprintable() or ch in "\n\t")
    alpha = sum(1 for ch in sample if ch.isalpha())
    ratio = printable / max(len(sample), 1)
    alpha_ratio = alpha / max(len(sample), 1)
    return ratio > 0.85 and alpha_ratio > 0.02


def extract_text_from_doc(path: str) -> str:
    """Extract text from legacy .doc files using system tools when available."""

    # macOS textutil can convert old Word docs to plain text
    if shutil.which("textutil"):
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", "-encoding", "UTF-8", path],
                check=True,
                capture_output=True,
                timeout=settings.doc_subprocess_extract_timeout_s,
            )
            text = result.stdout.decode("utf-8", errors="ignore")
            if _looks_like_text(text):
                return _normalize_whitespace_preserve_paragraphs(text)
            logger.warning("textutil output looked non-textual for %s; falling back", path)
        except subprocess.TimeoutExpired:
            logger.warning("textutil timed out extracting .doc from %s", path)
        except Exception as exc:
            logger.warning("textutil failed to extract .doc text from %s: %s", path, exc)
    # Fallback to antiword if installed
    if shutil.which("antiword"):
        try:
            result = subprocess.run(
                ["antiword", path],
                check=True,
                capture_output=True,
                timeout=settings.doc_subprocess_extract_timeout_s,
            )
            text = result.stdout.decode("utf-8", errors="ignore")
            if _looks_like_text(text):
                return _normalize_whitespace_preserve_paragraphs(text)
        except subprocess.TimeoutExpired:
            logger.warning("antiword timed out extracting .doc from %s", path)
        except Exception as exc:
            logger.warning("antiword failed to extract .doc text from %s: %s", path, exc)
    # Fallback to strings when other tools fail
    if shutil.which("strings"):
        try:
            result = subprocess.run(
                ["strings", "-a", path],
                check=True,
                capture_output=True,
                timeout=settings.doc_subprocess_extract_timeout_s,
            )
            text = result.stdout.decode("utf-8", errors="ignore")
            if _looks_like_text(text):
                return _normalize_whitespace_preserve_paragraphs(text)
        except subprocess.TimeoutExpired:
            logger.warning("strings timed out extracting .doc from %s", path)
        except Exception as exc:
            logger.warning("strings failed to extract .doc text from %s: %s", path, exc)
    return ""


def extract_text_from_docx(path: str) -> str:
    # If the file is not a DOCX zip, treat it as legacy DOC
    if not zipfile.is_zipfile(path):
        return extract_text_from_doc(path)
    # First try python-docx
    text = ""
    try:
        doc = Document(path)
        parts: List[str] = []
        for block in _iter_docx_blocks(doc):
            parts.append(block)
        text = "\n\n".join(p for p in parts if p)
        text = _normalize_whitespace_preserve_paragraphs(text)
        if not _looks_like_text(text):
            text = ""
    except Exception:
        text = ""

    # Fallback: parse document.xml directly
    if not text:
        try:
            with zipfile.ZipFile(path) as zf:
                with zf.open("word/document.xml") as fh:
                    data = fh.read()
            root = ET.fromstring(data)
            texts: List[str] = []
            for node in root.iter():
                if node.tag.endswith("}t") and node.text:
                    texts.append(node.text)
            xml_text = "\n".join(texts)
            xml_text = _normalize_whitespace_preserve_paragraphs(xml_text)
            if _looks_like_text(xml_text):
                text = xml_text
        except Exception:
            pass

    # macOS textutil can convert DOCX to text
    if not text and shutil.which("textutil"):
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", "-encoding", "UTF-8", path],
                check=True,
                capture_output=True,
                timeout=settings.doc_subprocess_extract_timeout_s,
            )
            txt = result.stdout.decode("utf-8", errors="ignore")
            if _looks_like_text(txt):
                text = _normalize_whitespace_preserve_paragraphs(txt)
        except subprocess.TimeoutExpired:
            logger.warning("textutil timed out extracting .docx from %s", path)
        except Exception as exc:
            logger.warning("textutil failed to extract .docx text from %s: %s", path, exc)

    # NOTE: avoid `strings` fallback for DOCX archives because it can produce
    # noisy ZIP/container artifacts that look text-like and explode chunk counts.
    # For DOCX, prefer structured extractors + embedded-image OCR only.

    embedded_ocr = _extract_embedded_image_ocr_from_zip(path, "word/media/")
    if embedded_ocr:
        text = "\n\n".join([part for part in [text, embedded_ocr] if part and part.strip()])
        text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def _recursive_split(text: str, chunk_size: int, separators: tuple[str, ...]) -> List[str]:
    if not text:
        return []
    if len(text) <= chunk_size or not separators:
        return [text]

    sep = separators[0]
    if sep:
        pieces = text.split(sep)
        rebuilt: List[str] = []
        buf = ""
        joiner = sep
        for piece in pieces:
            candidate = (buf + joiner + piece) if buf else piece
            if len(candidate) <= chunk_size:
                buf = candidate
            else:
                if buf:
                    rebuilt.append(buf)
                if len(piece) <= chunk_size:
                    buf = piece
                else:
                    rebuilt.extend(_recursive_split(piece, chunk_size, separators[1:]))
                    buf = ""
        if buf:
            rebuilt.append(buf)
        return rebuilt
    else:
        return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]


def _split_sentences(text: str) -> List[str]:
    if not text:
        return []
    splitter = (settings.sentence_splitter or "regex").lower()
    if splitter == "nltk":
        try:
            import nltk  # type: ignore

            try:
                nltk.data.find("tokenizers/punkt")
            except LookupError:
                nltk.download("punkt", quiet=True)
            return [s.strip() for s in nltk.sent_tokenize(text) if s and s.strip()]
        except Exception:
            splitter = "regex"
    if splitter == "spacy":
        try:
            import spacy  # type: ignore

            nlp = spacy.blank("en")
            if "sentencizer" not in nlp.pipe_names:
                nlp.add_pipe("sentencizer")
            doc = nlp(text)
            return [sent.text.strip() for sent in doc.sents if sent.text.strip()]
        except Exception:
            splitter = "regex"
    # Preserve punctuation by splitting on whitespace that follows sentence terminators.
    parts = SENTENCE_SPLIT_RE.split(text.strip())
    return [p.strip() for p in parts if p and p.strip()]


def _sentence_pack_chunk(text: str, chunk_size: int, separators: tuple[str, ...]) -> List[str]:
    """
    Pack chunks by paragraph -> sentence with hard fallback to recursive splitting
    when sentences exceed the chunk size.
    """
    if not text:
        return []
    paragraphs = [p.strip() for p in PARA_SPLIT_RE.split(text) if p.strip()]
    chunks: List[str] = []
    buf = ""

    def flush_buf() -> None:
        nonlocal buf
        if buf:
            chunks.append(buf.strip())
            buf = ""

    for para in paragraphs:
        sentences = _split_sentences(para) or [para]
        for sent in sentences:
            if len(sent) > chunk_size:
                flush_buf()
                chunks.extend(_recursive_split(sent, chunk_size, separators))
                continue
            candidate = f"{buf} {sent}".strip() if buf else sent
            if len(candidate) <= chunk_size:
                buf = candidate
            else:
                flush_buf()
                buf = sent
        flush_buf()
    return [c for c in chunks if c]


def _apply_overlap(chunks: List[str], overlap: int) -> List[str]:
    if overlap <= 0 or not chunks:
        return chunks
    out: List[str] = []
    prev_tail = ""
    for ch in chunks:
        prefix = prev_tail
        combined = f"{prefix}\n\n{ch}" if prefix else ch
        out.append(combined)
        prev_tail = ch[-overlap:]
    return out


def chunk_text(text: str, params: ChunkParams = ChunkParams()) -> List[str]:
    # Normalize while preserving paragraph boundaries; add extra spacing around likely headings
    text = _normalize_whitespace_preserve_paragraphs(text)
    text = _insert_heading_boundaries(text)

    if params.strategy == "sentence_pack":
        base_chunks = _sentence_pack_chunk(text, params.chunk_size, params.separators)
    else:
        base_chunks = _recursive_split(text, params.chunk_size, params.separators)
    if not base_chunks:
        return []
    return _apply_overlap(base_chunks, params.chunk_overlap)


def _tables_to_text(tables: Iterable[Iterable[Iterable[str | None]]]) -> str:
    """Convert table rows to a tab-separated text block."""
    rows: List[str] = []
    for table in tables:
        for row in table:
            if row is None:
                continue
            cells = [(cell or "").strip() for cell in row]
            if any(cells):
                rows.append("\t".join(cells))
    return "\n".join(rows)


def _iter_docx_blocks(doc: Document) -> List[str]:
    """Return paragraph and table text blocks in document order."""
    blocks: List[str] = []
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            para = Paragraph(child, doc)
            text = (para.text or "").strip()
            if text:
                blocks.append(text)
        elif isinstance(child, CT_Tbl):
            try:
                table = Table(child, doc)
                rows: List[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    blocks.append("\n".join(rows))
            except Exception:
                continue
    return blocks


def extract_text_from_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts: List[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts: List[str] = [f"Slide {slide_idx}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_parts.append(shape.text.strip())
            if hasattr(shape, "table") and shape.table is not None:
                table_rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        table_rows.append("\t".join(cells))
                if table_rows:
                    slide_parts.append("\n".join(table_rows))
        slide_text = "\n".join([p for p in slide_parts if p])
        if slide_text:
            parts.append(slide_text)
    text = "\n\n".join(parts)
    embedded_ocr = _extract_embedded_image_ocr_from_zip(path, "ppt/media/")
    if embedded_ocr:
        text = "\n\n".join([part for part in [text, embedded_ocr] if part and part.strip()])
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text


def extract_text_from_excel(path: str) -> str:
    import pandas as pd
    ext = os.path.splitext(path)[1].lower()
    engine = "openpyxl" if ext == ".xlsx" else "xlrd"
    parts: List[str] = []
    workbook = pd.ExcelFile(path, engine=engine)
    for sheet in workbook.sheet_names:
        df = workbook.parse(sheet)
        if df.empty:
            continue
        parts.append(f"Sheet: {sheet}")
        parts.append(df.to_csv(index=False, sep="\t"))
    text = "\n".join(parts)
    text = _normalize_whitespace_preserve_paragraphs(text)
    return text
